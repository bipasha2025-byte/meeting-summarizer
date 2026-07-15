import streamlit as st
import subprocess
import tempfile
import json
import os
import base64
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from groq import Groq

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Meeting Summarizer",
    page_icon="🎬",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Custom CSS ────────────────────────────────────────────────────────────────
st.markdown("""
<style>
  .main-header {
    background: #1f2328;
    padding: 2rem; border-radius: 12px; margin-bottom: 2rem; color: white;
  }
  .main-header h1 { margin: 0; font-size: 2rem; }
  .main-header p  { margin: 0.5rem 0 0; color: #aaa; font-size: 0.95rem; }
  .slide-preview {
    background: white; border: 1px solid #e5e7eb;
    border-radius: 8px; padding: 1rem; margin-bottom: 0.5rem;
    min-height: 120px;
  }
  .slide-title { font-weight: 700; color: #3b82d4; margin-bottom: 0.5rem; font-size: 0.9rem; }
</style>
""", unsafe_allow_html=True)

# ── Header ────────────────────────────────────────────────────────────────────
st.markdown("""
<div class="main-header">
  <h1>Meeting Summarizer</h1>
  <p>Upload a meeting video &rarr; AI transcribes &amp; analyses &rarr; Download PowerPoint</p>
</div>
""", unsafe_allow_html=True)

# ── Sidebar ───────────────────────────────────────────────────────────────────
with st.sidebar:
    st.header("Settings")

    default_key = ""
    try:
        default_key = st.secrets.get("GROQ_API_KEY", "")
    except Exception:
        pass

    api_key = st.text_input(
        "Groq API Key",
        value=default_key,
        type="password",
        help="Free key from console.groq.com — starts with gsk_",
        placeholder="gsk_...",
    )

    if api_key:
        if api_key.startswith("gsk_"):
            st.success("Key format looks correct (gsk_...)")
        else:
            st.error("Invalid key format. Groq keys must start with gsk_  — get yours free at console.groq.com")

    st.subheader("Meeting Info (optional)")
    meeting_title = st.text_input("Meeting Title", placeholder="e.g. Q4 Planning Session")
    meeting_date  = st.text_input("Meeting Date",  placeholder="e.g. 15 Jan 2025")

    st.divider()
    st.caption("Built for IBM WatsonX Challenge")
    st.caption("Powered by Groq (Llama 3.3 70B) + Whisper + python-pptx")
    st.caption("Free — no credit card needed")


# ── Helpers ───────────────────────────────────────────────────────────────────
GROQ_MAX_BYTES = 24 * 1024 * 1024
CHUNK_SECONDS  = 15 * 60


def _run_ffmpeg(cmd):
    try:
        result = subprocess.run(cmd, capture_output=True, text=True)
    except FileNotFoundError:
        raise RuntimeError(
            "ffmpeg is not installed on this server. If you're on Streamlit "
            "Community Cloud, add a line containing just 'ffmpeg' to "
            "packages.txt and redeploy."
        )
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg failed: {result.stderr[-800:]}")
    return result


def extract_compressed_audio(video_path, progress):
    progress.progress(10, "Extracting and compressing audio...")
    audio_path = f"{video_path}_audio.mp3"
    cmd = [
        "ffmpeg", "-y", "-i", video_path,
        "-vn", "-ac", "1", "-ar", "16000", "-b:a", "32k",
        audio_path,
    ]
    _run_ffmpeg(cmd)
    return audio_path


def split_audio_into_chunks(audio_path, chunk_seconds=CHUNK_SECONDS):
    base    = audio_path.rsplit(".", 1)[0]
    pattern = f"{base}_chunk_%03d.mp3"
    cmd = [
        "ffmpeg", "-y", "-i", audio_path,
        "-f", "segment", "-segment_time", str(chunk_seconds),
        "-c", "copy", pattern,
    ]
    _run_ffmpeg(cmd)
    glob_pattern = Path(pattern).name.replace("%03d", "*")
    chunk_files  = sorted(Path(audio_path).parent.glob(glob_pattern))
    return [str(p) for p in chunk_files]


def _transcribe_file(client, path):
    with open(path, "rb") as f:
        return client.audio.transcriptions.create(
            file=(Path(path).name, f),
            model="whisper-large-v3-turbo",
            response_format="text",
        )


def transcribe_video(video_path, key, progress):
    client     = Groq(api_key=key)
    audio_path = extract_compressed_audio(video_path, progress)
    audio_size = os.path.getsize(audio_path)

    try:
        if audio_size <= GROQ_MAX_BYTES:
            progress.progress(30, "Transcribing audio with Groq Whisper...")
            transcription = _transcribe_file(client, audio_path)
            progress.progress(50, "Transcription complete.")
            return transcription

        progress.progress(25, "Audio still large — splitting into chunks...")
        chunk_paths = split_audio_into_chunks(audio_path)
        texts = []
        n = len(chunk_paths)
        for i, chunk_path in enumerate(chunk_paths):
            pct = 30 + int((i / max(n, 1)) * 20)
            progress.progress(pct, f"Transcribing chunk {i + 1}/{n}...")
            texts.append(_transcribe_file(client, chunk_path))
            os.unlink(chunk_path)

        progress.progress(50, "Transcription complete.")
        return " ".join(texts)
    finally:
        if os.path.exists(audio_path):
            os.unlink(audio_path)


def summarize_with_groq(transcript, title, date, key, progress):
    """Extract key points from transcript using Groq Llama 3.3 70B."""
    progress.progress(55, "Analysing transcript with Llama 3.3 70B...")

    client = Groq(api_key=key)

    system_prompt = """You are an expert meeting analyst.
Extract the most important information from meeting transcripts.
Return ONLY valid JSON — no markdown fences, no extra text, just the raw JSON object.

Use this EXACT structure:
{
  "title": "meeting title",
  "date": "meeting date",
  "executive_summary": "2-3 sentence high-level summary",
  "key_points": ["key discussion point 1", "key discussion point 2"],
  "decisions": ["decision agreed upon 1"],
  "action_items": [
    {"task": "what needs to be done", "owner": "person or TBD", "due": "date or TBD"}
  ],
  "next_steps": ["follow-up item 1"],
  "attendees": ["Name 1", "Name 2"]
}

Rules:
- key_points must cover ALL major topics discussed
- decisions = only things explicitly agreed upon
- Extract action items even if implied; use TBD for unknown owners/dates
- Be concise but complete"""

    response = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Meeting Title: {title or 'Untitled Meeting'}\nMeeting Date: {date or 'Unknown'}\n\nTranscript:\n{transcript}"},
        ],
        temperature=0.2,
        response_format={"type": "json_object"},
    )

    summary = json.loads(response.choices[0].message.content)
    progress.progress(80, "Analysis complete.")
    return summary


def hex_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_title_slide(prs, title, date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_rgb("1f2328")

    txb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.3), Inches(1.5))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size  = Pt(36)
    run.font.bold  = True
    run.font.color.rgb = hex_rgb("ffffff")

    if date:
        txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5))
        p2   = txb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = date
        run2.font.size  = Pt(16)
        run2.font.color.rgb = hex_rgb("aaaaaa")


def add_content_slide(prs, title, items, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_rgb("ffffff")

    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.65))
    p   = txb.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size  = Pt(26)
    run.font.bold  = True
    run.font.color.rgb = hex_rgb(accent)

    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12.3), Emu(12700))
    ln.fill.background()
    ln.line.color.rgb = hex_rgb(accent)
    ln.line.width = Emu(25400)

    txb2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5))
    tf2  = txb2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(items):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        para.space_after = Pt(6)
        run2 = para.add_run()
        if isinstance(item, dict):
            run2.text = (
                f"  {item.get('task', '')}   |   "
                f"Owner: {item.get('owner', 'TBD')}   |   "
                f"Due: {item.get('due', 'TBD')}"
            )
        else:
            run2.text = f"  {item}"
        run2.font.size  = Pt(14)
        run2.font.color.rgb = hex_rgb("1f2328")


def generate_ppt(summary, progress):
    progress.progress(85, "Building PowerPoint slides...")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs, summary.get("title", "Meeting Summary"), summary.get("date", ""))

    if summary.get("executive_summary"):
        add_content_slide(prs, "Executive Summary", [summary["executive_summary"]], "3b82d4")

    if summary.get("key_points"):
        add_content_slide(prs, "Key Points", summary["key_points"], "3b82d4")

    if summary.get("decisions"):
        add_content_slide(prs, "Decisions Made", summary["decisions"], "7c5cd8")

    if summary.get("action_items"):
        add_content_slide(prs, "Action Items", summary["action_items"], "15803d")

    if summary.get("next_steps"):
        add_content_slide(prs, "Next Steps", summary["next_steps"], "c2410c")

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    data = Path(tmp.name).read_bytes()
    os.unlink(tmp.name)
    progress.progress(100, "Done!")
    return data


# ── Main UI ───────────────────────────────────────────────────────────────────
col1, col2 = st.columns([1, 1], gap="large")

with col1:
    st.subheader("1  Upload Meeting Video")
    uploaded = st.file_uploader(
        "Drop your video here",
        type=["mp4", "mov", "mkv", "avi", "webm", "m4a", "mp3", "wav"],
        help="Zoom, Teams, Meet, or any screen recording",
    )
    if uploaded:
        if uploaded.type.startswith("video"):
            st.video(uploaded)
        st.caption(f"{uploaded.name}  ({uploaded.size / 1024 / 1024:.1f} MB)")

with col2:
    st.subheader("2  Generate Summary + PPT")

    if not api_key:
        st.info("Enter your Groq API key in the sidebar. Get one free at console.groq.com — starts with gsk_")
    elif not api_key.startswith("gsk_"):
        st.error("Your key must start with gsk_  — go to console.groq.com → API Keys → Create API Key")
    elif not uploaded:
        st.info("Upload a meeting video on the left to begin.")
    else:
        if st.button("Generate Summary + PPT", type="primary", use_container_width=True):
            progress = st.progress(0, "Starting...")
            try:
                progress.progress(5, "Saving uploaded file...")
                suffix    = Path(uploaded.name).suffix.lower() or ".mp4"
                tmp_video = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
                tmp_video.write(uploaded.getbuffer())
                tmp_video.close()

                # Step 1: Transcribe with Groq Whisper
                transcript = transcribe_video(tmp_video.name, api_key, progress)
                st.session_state["transcript"] = transcript

                # Step 2: Summarise with Llama 3.3
                summary = summarize_with_groq(
                    transcript,
                    meeting_title or uploaded.name,
                    meeting_date or "",
                    api_key,
                    progress,
                )
                st.session_state["summary"] = summary

                # Step 3: Build PPT
                ppt_bytes = generate_ppt(summary, progress)
                st.session_state["ppt_bytes"] = ppt_bytes
                os.unlink(tmp_video.name)

            except Exception as e:
                progress.empty()
                st.error(f"Error: {e}")
                st.stop()

        if "summary" in st.session_state:
            summary   = st.session_state["summary"]
            ppt_bytes = st.session_state["ppt_bytes"]
            fname = (summary.get("title") or "meeting_summary").replace(" ", "_") + ".pptx"

            st.success("Processing complete!")
            st.download_button(
                label="Download PowerPoint (.pptx)",
                data=ppt_bytes,
                file_name=fname,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary",
            )

            if st.session_state.get("transcript"):
                with st.expander("View Transcript"):
                    st.text_area(
                        "Transcript",
                        st.session_state["transcript"],
                        height=200,
                        label_visibility="collapsed",
                    )

# ── Slide Preview ─────────────────────────────────────────────────────────────
if "summary" in st.session_state:
    s = st.session_state["summary"]
    st.divider()
    st.subheader("Slide Preview")

    slides = [
        ("Slide 1 — Title",            [s.get("title", ""), s.get("date", "")]),
        ("Slide 2 — Executive Summary", [s.get("executive_summary", "")]),
        ("Slide 3 — Key Points",        s.get("key_points", [])),
        ("Slide 4 — Decisions Made",    s.get("decisions", [])),
        ("Slide 5 — Action Items",
         [f"{a.get('task','')} | {a.get('owner','TBD')} | {a.get('due','TBD')}"
          for a in s.get("action_items", [])]),
        ("Slide 6 — Next Steps",        s.get("next_steps", [])),
    ]

    cols = st.columns(3)
    for i, (slide_title, items) in enumerate(slides):
        with cols[i % 3]:
            content = "".join(
                f"<div style='font-size:12px;margin:2px 0'>• {it}</div>"
                for it in items if it
            )
            st.markdown(
                f'<div class="slide-preview">'
                f'<div class="slide-title">{slide_title}</div>{content}</div>',
                unsafe_allow_html=True,
            )
